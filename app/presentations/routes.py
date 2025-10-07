import os
import glob
from flask import render_template, request, redirect, url_for, flash, jsonify, send_from_directory, session
from flask_login import login_required, current_user
from werkzeug.utils import secure_filename
from . import presentations_bp
from .models import Presentation, Slide, MediaFile
from ..extensions import db
from pptx import Presentation as PptxPresentation
import shutil
import re
from pptx.enum.shapes import MSO_SHAPE_TYPE
from ..utils.thumbnails import auto_generate_thumbnail

import zipfile, json, shutil
from flask import send_file
from ..teams.models import Team

from pdf2image import convert_from_path
from setup_poppler import ensure_poppler
from flask import send_from_directory, current_app



UPLOAD_PPT_DIR = os.path.join(os.getcwd(), "uploads", "ppts")
UPLOAD_MEDIA_DIR = os.path.join(os.getcwd(), "uploads", "media")
UPLOAD_MEDIA_PPT_DIR = os.path.join(os.getcwd(), "uploads", "media", "slides")
UPLOAD_MEDIA_SLIDES_DIR = os.path.join(os.getcwd(), "uploads", "media", "ppts")
os.makedirs(UPLOAD_PPT_DIR, exist_ok=True)
os.makedirs(UPLOAD_MEDIA_DIR, exist_ok=True)
os.makedirs(UPLOAD_MEDIA_PPT_DIR, exist_ok=True)
os.makedirs(UPLOAD_MEDIA_SLIDES_DIR, exist_ok=True)


def import_pptx(file_path, pres_id):
    ppt = PptxPresentation(file_path)
    idx = 0
    for i in range(0, len(ppt.slides), 2):
        client_html = extract_slide_content(ppt.slides[i])

        notes_html = ""
        if i + 1 < len(ppt.slides):
            notes_html = extract_slide_content(ppt.slides[i+1])

        new_slide = Slide(
            presentation_id=pres_id,
            position=idx,
            client_content=client_html,
            agent_notes=notes_html
        )
        db.session.add(new_slide)
        idx += 1
    db.session.commit()

def import_pdf_as_images(pdf_path, pres_id):
    poppler_path = ensure_poppler()  # Windows returns path, Linux/mac just None
    pages = convert_from_path(pdf_path, dpi=150, poppler_path=poppler_path)
    for idx, page in enumerate(pages):
        fname = f"slide_{pres_id}_{idx}.jpg"
        fpath = os.path.join(UPLOAD_MEDIA_DIR, "slides", fname)
        page.save(fpath, "JPEG")
        media = MediaFile(filename="slides/"+fname, user_id=current_user.id)
        db.session.add(media)
        slide = Slide(presentation_id=pres_id, position=idx,
                      client_content=f'<img src="{url_for("presentations.media_file", "slides", filename=fname)}">', agent_notes="")
        db.session.add(slide)


    db.session.commit()


@presentations_bp.route("/")
@login_required
def index():
    pres = Presentation.query.all()
    return render_template("presentations/index.html", presentations=pres)

@presentations_bp.route("/create", methods=["GET","POST"])
@login_required
def create():
    if current_user.role not in ("MANAGER","ADMIN","SUPER_ADMIN"):
        flash("Denied","danger"); return redirect(url_for("presentations.index"))
    if request.method=="POST":
        title=request.form["title"]
        pres=Presentation(title=title,creator_id=current_user.id)
        db.session.add(pres); db.session.commit()

        # Option A: manual tinymce slides
        slides=request.form.getlist("client_content"); notes=request.form.getlist("agent_notes")
        for idx,(c,a) in enumerate(zip(slides,notes)):
            s=Slide(presentation_id=pres.id,position=idx,client_content=c,agent_notes=a)
            db.session.add(s)
        db.session.commit()

        # Option B: PPT uploaded
        ppt_file=request.files.get("pptx")
        if ppt_file and ppt_file.filename.endswith(".pptx"):
            fname=secure_filename(ppt_file.filename)
            fpath=os.path.join(UPLOAD_PPT_DIR,fname)
            ppt_file.save(fpath)
            import_pptx(fpath, pres.id)

        flash("Presentation created","success")
        return redirect(url_for("presentations.index"))
    return render_template("presentations/create.html")

@presentations_bp.route("/<int:pres_id>/edit", methods=["GET","POST"])
@login_required
def edit(pres_id):
    pres=Presentation.query.get_or_404(pres_id)
    if current_user.role not in ("MANAGER","ADMIN","SUPER_ADMIN"):
        flash("Denied","danger"); return redirect(url_for("presentations.index"))

    if request.method=="POST":
        pres.title=request.form["title"]
        # Clear slides then rebuild
        Slide.query.filter_by(presentation_id=pres.id).delete()
        slides=request.form.getlist("client_content"); notes=request.form.getlist("agent_notes")
        for idx,(c,a) in enumerate(zip(slides,notes)):
            s=Slide(presentation_id=pres.id,position=idx,client_content=c,agent_notes=a)
            db.session.add(s)
        db.session.commit()
        # Re-import if PPT uploaded
        ppt_file=request.files.get("pptx")
        if ppt_file and ppt_file.filename.endswith(".pptx"):
            fname=secure_filename(ppt_file.filename)
            fpath=os.path.join(UPLOAD_PPT_DIR,fname)
            ppt_file.save(fpath)
            import_pptx(fpath, pres.id)
        db.session.commit()
        flash("Presentation updated","success")
        return redirect(url_for("presentations.index"))
    return render_template("presentations/edit.html", presentation=pres)

@presentations_bp.route("/<int:pres_id>/run_agent")
@login_required
def run_agent(pres_id):
    pres=Presentation.query.get_or_404(pres_id)
    return render_template("presentations/run_agent.html", presentation=pres)

@presentations_bp.route("/<int:pres_id>/run_client")
def run_client(pres_id):
    pres=Presentation.query.get_or_404(pres_id)
    return render_template("presentations/run_client.html", presentation=pres)

current_slide={}
@presentations_bp.route("/<int:pres_id>/goto/<int:slide_id>", methods=["POST"])
@login_required
def goto_slide(pres_id,slide_id):
    current_slide[pres_id]=slide_id
    return jsonify({"ok":True,"slide":slide_id})

@presentations_bp.route("/<int:pres_id>/current")
def current(pres_id):
    sid=current_slide.get(pres_id)
    if not sid: return jsonify({"slide":None})
    slide=Slide.query.get(sid)
    return jsonify({"id":slide.id,"client_content":slide.client_content,"agent_notes":slide.agent_notes})

@presentations_bp.route("/upload_media", methods=["POST"])
@login_required
def upload_media():
    if current_user.role not in ("MANAGER","ADMIN","SUPER_ADMIN"):
        return jsonify({"error": "Unauthorized"}), 403

    file = request.files.get("file")
    if not file:
        flash("No file selected", "danger")
        return redirect(url_for("presentations.media_manager"))

    fname = f"{current_user.id}_{secure_filename(file.filename)}"
    save_dir = os.path.join("uploads", "media")
    os.makedirs(save_dir, exist_ok=True)
    file_path = os.path.join(save_dir, fname)
    file.save(file_path)

    # âœ… Auto-generate thumbnail (video, image, pdf)
    thumb_name = auto_generate_thumbnail(file_path, file.mimetype, save_dir)

    # Write DB record
    media = MediaFile(filename=fname, user_id=current_user.id)
    db.session.add(media)
    db.session.commit()

    # If upload came from a form, redirect instead of JSON
    if request.form:
        flash("Media uploaded", "success")
        return redirect(url_for("presentations.media_manager"))

    file_url = url_for("presentations.media_file", filename=fname, _external=True)
    return jsonify({"location": file_url, "thumbnail": thumb_name})

# Serve uploaded media
@presentations_bp.route("/media/<path:filename>")
def media_file(filename):
    """Serve uploaded media files."""
    media_dir = os.path.join(os.getcwd(), "uploads", "media")
    return send_from_directory(media_dir, filename)

@presentations_bp.route("/<int:pres_id>/duplicate", methods=["POST"])
@login_required
def duplicate(pres_id):
    pres = Presentation.query.get_or_404(pres_id)
    if current_user.role not in ("MANAGER","ADMIN","SUPER_ADMIN"):
        flash("Not authorized", "danger")
        return redirect(url_for("presentations.index"))

    # Create copy
    new_pres = Presentation(title=pres.title + " (Copy)", creator_id=current_user.id)
    db.session.add(new_pres)
    db.session.flush()

    for slide in pres.slides:
        new_slide = Slide(
            presentation_id=new_pres.id,
            position=slide.position,
            client_content=slide.client_content,
            agent_notes=slide.agent_notes,
        )
        db.session.add(new_slide)

    db.session.commit()
    flash("Presentation duplicated", "success")
    return redirect(url_for("presentations.edit", pres_id=new_pres.id))


@presentations_bp.route("/media_manager")
@login_required
def media_manager():
    if current_user.role in ("ADMIN","SUPER_ADMIN"):
        files = MediaFile.query.all()
    elif current_user.role == "MANAGER":
        files = MediaFile.query.filter_by(user_id=current_user.id).all()
    else:
        flash("Unauthorized", "danger")
        return redirect(url_for("presentations.index"))
    
    return render_template("presentations/media_manager.html", files=files)

@presentations_bp.route("/media_manager/delete/<filename>", methods=["POST"])
@login_required
def delete_media(filename):
    if current_user.role not in ("ADMIN","SUPER_ADMIN","MANAGER"):
        return jsonify({"ok":False,"error":"Unauthorized"}),403

    media_dir = os.path.join("uploads","media")
    fpath = os.path.join(media_dir, filename)

    # âœ… Delete physical file if present
    if os.path.exists(fpath):
        os.remove(fpath)

    # âœ… Remove auto-generated thumbnail (_thumb.jpg)
    base, _ = os.path.splitext(filename)
    thumb = base + "_thumb.jpg"
    thumb_path = os.path.join(media_dir, thumb)
    if os.path.exists(thumb_path):
        os.remove(thumb_path)

    # âœ… ALSO delete DB record
    from .models import MediaFile
    media = MediaFile.query.filter_by(filename=filename).first()
    if media:
        db.session.delete(media)
        db.session.commit()

    return jsonify({"ok":True})


@presentations_bp.route("/import_ppt", methods=["GET","POST"])
@login_required
def import_ppt():
    if current_user.role not in ("MANAGER","ADMIN","SUPER_ADMIN"):
        flash("Unauthorized","danger")
        return redirect(url_for("presentations.index"))

    if request.method == "POST":
        title = request.form["title"]
        restricted = bool(request.form.get("restricted_to_managers"))
        reverse_order = "reverse" in request.form
        skip_notes = "skip_notes" in request.form
        ppt_file = request.files.get("pptx")

        if not ppt_file or not (ppt_file.filename.endswith(".pptx") or ppt_file.filename.endswith(".ppt")):
            flash("Please upload a .pptx or .ppt file","danger")
            return redirect(url_for("presentations.import_ppt"))

        fname = secure_filename(ppt_file.filename)
        fpath = os.path.join(UPLOAD_PPT_DIR, fname)
        ppt_file.save(fpath)

        # âœ… If old .ppt, convert to .pptx with LibreOffice
        if fname.endswith(".ppt"):
            converted = os.path.splitext(fpath)[0] + ".pptx"
            os.system(f'libreoffice --headless --convert-to pptx "{fpath}" --outdir "{UPLOAD_PPT_DIR}"')
            fpath = converted

        # Create new Presentation record
        pres = Presentation(
            title=title,
            creator_id=current_user.id,
            restricted_to_managers=restricted
        )
        db.session.add(pres)
        db.session.commit()

        # Parse pptx
        ppt = PptxPresentation(fpath)
        slides = ppt.slides

        slide_pairs = []
        for i in range(0, len(slides), 2):
            client_slide = slides[i]
            notes_slide = slides[i+1] if i+1 < len(slides) else None
            slide_pairs.append((client_slide, notes_slide))

        if reverse_order:
            slide_pairs = [(n, c) for c, n in slide_pairs]

        for idx, (client, notes) in enumerate(slide_pairs):
            client_html = extract_slide_content(client)
            notes_html = "" if (skip_notes or not notes) else extract_slide_content(notes)

            s = Slide(
                presentation_id=pres.id,
                position=idx,
                client_content=client_html,
                agent_notes=notes_html
            )
            db.session.add(s)

        db.session.commit()
        flash("Presentation imported successfully","success")
        return redirect(url_for("presentations.index"))

    return render_template("presentations/import_ppt.html")


@presentations_bp.route("/confirm_import", methods=["POST"])
@login_required
def confirm_import():
    if current_user.role not in ("MANAGER","ADMIN","SUPER_ADMIN"):
        flash("Unauthorized","danger")
        return redirect(url_for("presentations.index"))

    slides_preview = session.get("ppt_preview")
    title = session.get("ppt_title")
    if not slides_preview:
        flash("No import data in session","danger")
        return redirect(url_for("presentations.import_ppt"))

    pres = Presentation(title=title, creator_id=current_user.id)
    db.session.add(pres)
    db.session.commit()

    for idx,slide in enumerate(slides_preview):
        s = Slide(presentation_id=pres.id, position=idx,
                  client_content=slide["client"], agent_notes=slide["notes"])
        db.session.add(s)

    db.session.commit()
    session.pop("ppt_preview", None)
    session.pop("ppt_title", None)

    flash("Presentation saved!","success")
    return redirect(url_for("presentations.index"))




def extract_slide_content(slide):
    """
    Extract text/media from pptx slide into semantic HTML.
    Handles: titles, paragraphs, bullet/numbered lists, images, video, audio.
    """

    content = []

    for shape in slide.shapes:
        # --- Titles & subtitles ---
        if shape.is_placeholder and shape.has_text_frame and shape.text.strip():
            if shape.placeholder_format.type == 1:  # Title
                content.append(f"<h1>{shape.text.strip()}</h1>")
                continue
            elif shape.placeholder_format.type == 3:  # Subtitle
                content.append(f"<h2>{shape.text.strip()}</h2>")
                continue

        # --- Text frames (paragraphs, bullet lists) ---
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                txt = p.text.strip()
                if not txt:
                    continue

                # Indented bullets
                indent = f" style='margin-left:{p.level*2}em'" if p.level > 0 else ""

                # Numbered lists (starts with "1.", "2.", etc.)
                if txt[:1].isdigit() and "." in txt[:3]:
                    parts = txt.split(".", 1)
                    txt_clean = parts[1].strip() if len(parts) == 2 else txt
                    content.append(f"<ol{indent}><li>{txt_clean}</li></ol>")
                # Bullets (â€¢ or -)
                elif txt.startswith("â€¢") or txt.startswith("- "):
                    txt_clean = txt.lstrip("â€¢- ").strip()
                    content.append(f"<ul{indent}><li>{txt_clean}</li></ul>")
                else:
                    content.append(f"<p{indent}>{txt}</p>")

        # --- Images ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
            try:
                image = shape.image
                ext = image.ext  # usually 'png', 'jpeg'
                img_name = f"{current_user.id}_{image.sha1}.{ext}"
                fpath = os.path.join(UPLOAD_MEDIA_DIR, img_name)

                if not os.path.exists(fpath):
                    with open(fpath, "wb") as f:
                        f.write(image.blob)

                url = url_for("presentations.media_file", filename=img_name, _external=False)
                content.append(f'<img src="{url}" class="slide-img" />')
            except Exception as e:
                print(f"Image extract error: {e}")

        # --- Embedded media (audio/video placeholders) ---
        if hasattr(shape, "media_type"):
            if shape.media_type == "video":
                content.append('<video controls class="slide-video"><source src="movie.mp4" type="video/mp4"></video>')
            elif shape.media_type == "audio":
                content.append('<audio controls class="slide-audio"><source src="audio.mp3" type="audio/mpeg"></audio>')

    return "\n".join(content) if content else "<p>(empty slide)</p>"



@presentations_bp.route("/<int:pres_id>/delete", methods=["POST"])
@login_required
def delete_presentation(pres_id):
    pres = Presentation.query.get_or_404(pres_id)
    if current_user.role not in ("MANAGER", "ADMIN", "SUPER_ADMIN"):
        flash("Not authorized", "danger")
        return redirect(url_for("presentations.index"))

    # Delete slides first
    Slide.query.filter_by(presentation_id=pres.id).delete()
    db.session.delete(pres)
    db.session.commit()
    flash("Presentation deleted", "success")
    return redirect(url_for("presentations.index"))



@presentations_bp.route("/search")
@login_required
def search():
    q = request.args.get("q", "").strip()
    results = []
    if q:
        results = Presentation.query.join(Slide).filter(
            (Presentation.title.ilike(f"%{q}%")) |
            (Slide.client_content.ilike(f"%{q}%")) |
            (Slide.agent_notes.ilike(f"%{q}%"))
        ).all()
    return render_template("presentations/search.html", results=results, q=q)



def can_view_presentation(pres, user):
    if pres.restricted_to_managers and user.role == "AGENT":
        return False
    if pres.access_rules:
        allowed = False
        if any(ar.user_id == user.id for ar in pres.access_rules):
            allowed = True
        team_ids = [tm.team_id for tm in user.team_memberships]
        if any(ar.team_id in team_ids for ar in pres.access_rules if ar.team_id):
            allowed = True
        if not allowed and user.role not in ("ADMIN", "SUPER_ADMIN"):
            return False
    return True

def import_zpb(zip_path):
    import zipfile, json
    with zipfile.ZipFile(zip_path, 'r') as z:
        details = json.loads(z.read("details.json").decode())
        pres = Presentation(title=details["title"], creator_id=current_user.id,
                            restricted_to_managers=details.get("restricted_to_managers", False))
        db.session.add(pres); db.session.commit()

        # import slides
        client_files = sorted([f for f in z.namelist() if f.startswith("client/")])
        agent_files = sorted([f for f in z.namelist() if f.startswith("agent/")])
        for idx, c in enumerate(client_files):
            client_html = z.read(c).decode()
            notes_html = ""
            if idx < len(agent_files):
                notes_html = z.read(agent_files[idx]).decode()
            db.session.add(Slide(presentation_id=pres.id, position=idx,
                                 client_content=client_html, agent_notes=notes_html))
        db.session.commit()

import tempfile

@presentations_bp.route("/<int:pres_id>/export_zpb")
@login_required
def export_zpb(pres_id):
    pres = Presentation.query.get_or_404(pres_id)
    if not can_view_presentation(pres, current_user):
        flash("Access denied", "danger")
        return redirect(url_for("presentations.index"))

    # Use system temp dir for portability
    tmpdir = os.path.join(tempfile.gettempdir(), f"pres_{pres.id}")
    if os.path.exists(tmpdir):
        shutil.rmtree(tmpdir)
    os.makedirs(os.path.join(tmpdir, "client"), exist_ok=True)
    os.makedirs(os.path.join(tmpdir, "agent"), exist_ok=True)
    os.makedirs(os.path.join(tmpdir, "media"), exist_ok=True)

    # dump slides as html
    for idx, s in enumerate(sorted(pres.slides, key=lambda x: x.position)):
        with open(os.path.join(tmpdir, "client", f"slide{idx}.html"), "w", encoding="utf-8") as f:
            f.write(s.client_content or "")
        with open(os.path.join(tmpdir, "agent", f"notes{idx}.html"), "w", encoding="utf-8") as f:
            f.write(s.agent_notes or "")

    # details.json
    details = {
        "id": pres.id,
        "title": pres.title,
        "restricted_to_managers": pres.restricted_to_managers,
        "teams": [ar.team_id for ar in pres.access_rules if ar.team_id],
        "users": [ar.user_id for ar in pres.access_rules if ar.user_id],
        "type": "html",
    }
    with open(os.path.join(tmpdir, "details.json"), "w", encoding="utf-8") as f:
        json.dump(details, f)

    # copy referenced media
    upload_dir = os.path.join("uploads", "media")
    for slide in pres.slides:
        html = (slide.client_content or "") + (slide.agent_notes or "")
        for fname in re.findall(r'/media/([^"\'\s>]+)', html):
            # normalize to OS path
            src = os.path.join(upload_dir, *fname.split("/"))
            dest = os.path.join(tmpdir, "media", *fname.split("/"))
            if os.path.exists(src):
                os.makedirs(os.path.dirname(dest), exist_ok=True)
                shutil.copy(src, dest)

    # zip it
    zpath = os.path.join(tempfile.gettempdir(), f"pres_{pres.id}.zpb")
    if os.path.exists(zpath):
        os.remove(zpath)
    with zipfile.ZipFile(zpath, "w") as z:
        for root, dirs, files in os.walk(tmpdir):
            for file in files:
                fpath = os.path.join(root, file)
                z.write(fpath, os.path.relpath(fpath, tmpdir))

    return send_file(zpath, as_attachment=True, download_name=f"{pres.title}.zpb")



@presentations_bp.route("/import_zpb", methods=["GET","POST"])
@login_required
def import_zpb():
    if request.method=="POST":
        f=request.files["file"]
        savepath=f"/tmp/{secure_filename(f.filename)}"
        f.save(savepath)
        with zipfile.ZipFile(savepath,"r") as z:
            details=json.loads(z.read("details.json").decode())
            pres=Presentation(title=details["title"], creator_id=current_user.id,
                              restricted_to_managers=details.get("restricted_to_managers",False))
            db.session.add(pres); db.session.commit()

            # slides
            client=[x for x in z.namelist() if x.startswith("client/")]
            agent=[x for x in z.namelist() if x.startswith("agent/")]
            for idx,c in enumerate(sorted(client)):
                client_html=z.read(c).decode()
                notes_html=z.read(agent[idx]).decode() if idx<len(agent) else ""
                db.session.add(Slide(presentation_id=pres.id,position=idx,
                                     client_content=client_html,agent_notes=notes_html))
            db.session.commit()

            # restrictions
            for tid in details.get("teams",[]): db.session.add(PresentationAccess(presentation_id=pres.id, team_id=tid))
            for uid in details.get("users",[]): db.session.add(PresentationAccess(presentation_id=pres.id, user_id=uid))
            db.session.commit()

            # media copy
            upload_dir=os.path.join("uploads","media"); os.makedirs(upload_dir,exist_ok=True)
            for f in z.namelist():
                if f.startswith("media/"):
                    dest=os.path.join(upload_dir,os.path.basename(f))
                    with open(dest,"wb") as out: out.write(z.read(f))
        flash("Backup imported","success"); return redirect(url_for("presentations.index"))
    return render_template("presentations/import_zpb.html")


@presentations_bp.route("/export_all_zpb")
@login_required
def export_all_zpb():
    if current_user.role not in ("SUPER_ADMIN", "ADMIN"):
        flash("Only Admins/Super Admins can export full backups", "danger")
        return redirect(url_for("presentations.index"))

    base_tmp = f"/tmp/presentations_all"
    if os.path.exists(base_tmp):
        shutil.rmtree(base_tmp)
    os.makedirs(base_tmp, exist_ok=True)

    all_pres = Presentation.query.all()
    upload_dir = os.path.join("uploads", "media")

    manifest = []  # ðŸ”¥ collect info for manifest.json

    for pres in all_pres:
        pres_dir = os.path.join(base_tmp, f"presentation_{pres.id}")
        os.makedirs(os.path.join(pres_dir, "client"), exist_ok=True)
        os.makedirs(os.path.join(pres_dir, "agent"), exist_ok=True)
        os.makedirs(os.path.join(pres_dir, "media"), exist_ok=True)

        def rewrite_and_copy(html):
            if not html:
                return ""
            updated_html = html
            for match in re.findall(r'src="[^"]+/(uploads/media|presentations/media)/([^"]+)"', html):
                fname = match[1]
                src = os.path.join(upload_dir,  *fname.split("/"))
                dest = os.path.join(pres_dir, "media",  *fname.split("/"))
                if os.path.exists(src):
                    os.makedirs(os.path.dirname(dest), exist_ok=True)
                    shutil.copy(src, dest)
                    updated_html = updated_html.replace(f"/presentations/media/{fname}", f"media/{fname}")
                    updated_html = updated_html.replace(f"/uploads/media/{fname}", f"media/{fname}")
            return updated_html

        for idx, s in enumerate(sorted(pres.slides, key=lambda x: x.position)):
            with open(os.path.join(pres_dir, "client", f"slide{idx}.html"), "w") as f:
                f.write(rewrite_and_copy(s.client_content or ""))
            with open(os.path.join(pres_dir, "agent", f"notes{idx}.html"), "w") as f:
                f.write(rewrite_and_copy(s.agent_notes or ""))

        details = {
            "id": pres.id,
            "title": pres.title,
            "restricted_to_managers": pres.restricted_to_managers,
            "teams": [ar.team_id for ar in pres.access_rules if ar.team_id],
            "users": [ar.user_id for ar in pres.access_rules if ar.user_id],
            "type": "html"
        }
        with open(os.path.join(pres_dir, "details.json"), "w") as f:
            json.dump(details, f)

        # add entry to manifest
        manifest.append({
            "id": pres.id,
            "title": pres.title,
            "created_at": pres.created_at.strftime("%Y-%m-%d %H:%M"),
            "restricted_to_managers": pres.restricted_to_managers,
            "teams": details["teams"],
            "users": details["users"],
            "slide_count": len(pres.slides)
        })

    # write manifest.json at root
    with open(os.path.join(base_tmp, "manifest.json"), "w") as f:
        json.dump({"presentations": manifest}, f, indent=2)

    zpath = "/tmp/all_presentations_backup.zip"
    if os.path.exists(zpath):
        os.remove(zpath)

    with zipfile.ZipFile(zpath, "w") as z:
        for root, dirs, files in os.walk(base_tmp):
            for file in files:
                fpath = os.path.join(root, file)
                z.write(fpath, os.path.relpath(fpath, base_tmp))

    return send_file(
        zpath,
        as_attachment=True,
        download_name="all_presentations_backup.zip"
    )


@presentations_bp.route("/import_all_zpb", methods=["GET","POST"])
@login_required
def import_all_zpb():
    if current_user.role not in ("SUPER_ADMIN","ADMIN"):
        flash("Only Admins/Super Admins can import backups", "danger")
        return redirect(url_for("presentations.index"))

    if request.method == "POST":
        file = request.files["file"]
        savepath = f"/tmp/{secure_filename(file.filename)}"
        file.save(savepath)

        with zipfile.ZipFile(savepath, "r") as z:
            manifest = {}
            if "manifest.json" in z.namelist():
                manifest = json.loads(z.read("manifest.json").decode())
            presentations = manifest.get("presentations", [])

        # Step 1: show preview selection of which to import
        return render_template("presentations/import_preview_select.html",
                               file_path=savepath, presentations=presentations)

    return render_template("presentations/import_all_zpb.html")


@presentations_bp.route("/import_selected_process", methods=["POST"])
@login_required
def import_selected_process():
    if current_user.role not in ("SUPER_ADMIN","ADMIN"):
        flash("Unauthorized", "danger")
        return redirect(url_for("presentations.index"))

    savepath = request.form["file_path"]
    selected_ids = list(map(int, request.form.getlist("presentation_ids")))
    print(selected_ids)

    with zipfile.ZipFile(savepath, "r") as z:
        pres_folders = set([name.split("/")[0] for name in z.namelist() if name.startswith("presentation_")])
        upload_dir = os.path.join("uploads", "media")
        os.makedirs(upload_dir, exist_ok=True)

        restored = []
        for folder in pres_folders:
            details = json.loads(z.read(f"{folder}/details.json").decode())
            if details.get("id") not in selected_ids and details["title"] not in selected_ids:
                continue

            pres = Presentation(
                title=details["title"],
                creator_id=current_user.id,
                restricted_to_managers=details.get("restricted_to_managers", False)
            )
            db.session.add(pres); db.session.commit()

            client_files = sorted([x for x in z.namelist() if x.startswith(f"{folder}/client/")])
            agent_files = sorted([x for x in z.namelist() if x.startswith(f"{folder}/agent/")])

            def restore_media_refs(html):
                if not html: return ""
                updated_html = html
                for fname in re.findall(r'media/([^"^\']+)', html):
                    if f"{folder}/media/{fname}" in z.namelist():
                        dest = os.path.join(upload_dir, fname)
                        with open(dest,"wb") as out: out.write(z.read(f"{folder}/media/{fname}"))
                    updated_html = re.sub(
                        rf'(?<!presentations/)media/{re.escape(fname)}',
                        f'/presentations/media/{fname}',
                        updated_html
                    )
                return updated_html







            for idx, c in enumerate(client_files):
                client_html = restore_media_refs(z.read(c).decode())
                notes_html = restore_media_refs(z.read(agent_files[idx]).decode()) if idx < len(agent_files) else ""
                db.session.add(Slide(presentation_id=pres.id, position=idx,
                                     client_content=client_html, agent_notes=notes_html))
            db.session.commit()

            for tid in details.get("teams", []):
                db.session.add(PresentationAccess(presentation_id=pres.id, team_id=tid))
            for uid in details.get("users", []):
                db.session.add(PresentationAccess(presentation_id=pres.id, user_id=uid))
            db.session.commit()

            restored.append(pres.title)

    flash(f"Imported {len(restored)} presentations: {', '.join(restored)}", "success")
    return redirect(url_for("presentations.index"))






@presentations_bp.route("/export_select", methods=["GET","POST"])
@login_required
def export_select_zpb():
    if current_user.role not in ("SUPER_ADMIN","ADMIN"):
        flash("Unauthorized", "danger")
        return redirect(url_for("presentations.index"))

    if request.method == "POST":
        selected_ids = request.form.getlist("presentation_ids")
        if not selected_ids:
            flash("No presentations selected!", "warning")
            return redirect(url_for("presentations.export_select_zpb"))

        # reuse export logic but only for selected presentations
        import tempfile
        base_tmp = base_tmp = tempfile.mkdtemp(prefix="export_selected_")
        if os.path.exists(base_tmp):
            shutil.rmtree(base_tmp)
        os.makedirs(base_tmp, exist_ok=True)

        selected = Presentation.query.filter(Presentation.id.in_(selected_ids)).all()
        manifest = []
        upload_dir = os.path.join("uploads", "media")

        for pres in selected:
            pres_dir = os.path.join(base_tmp, f"presentation_{pres.id}")
            os.makedirs(os.path.join(pres_dir, "client"), exist_ok=True)
            os.makedirs(os.path.join(pres_dir, "agent"), exist_ok=True)
            os.makedirs(os.path.join(pres_dir, "media"), exist_ok=True)

            def rewrite_and_copy(html):
                if not html: return ""
                updated_html = html
                for match in re.findall(r'src="[^"]+/(uploads/media|presentations/media)/([^"]+)"', html):
                    fname = match[1]
                    src = os.path.join(upload_dir,  *fname.split("/"))
                    dest = os.path.join(pres_dir, "media",  *fname.split("/"))
                    if os.path.exists(src):
                        os.makedirs(os.path.dirname(dest), exist_ok=True)
                        shutil.copy(src, dest)
                        updated_html = updated_html.replace(f"/presentations/media/{fname}", f"media/{fname}")
                        updated_html = updated_html.replace(f"/uploads/media/{fname}", f"media/{fname}")
                return updated_html

            for idx, s in enumerate(sorted(pres.slides, key=lambda x: x.position)):
                with open(os.path.join(pres_dir, "client", f"slide{idx}.html"), "w") as f:
                    f.write(rewrite_and_copy(s.client_content or ""))
                with open(os.path.join(pres_dir, "agent", f"notes{idx}.html"), "w") as f:
                    f.write(rewrite_and_copy(s.agent_notes or ""))

            details = {
                "id": pres.id,
                "title": pres.title,
                "restricted_to_managers": pres.restricted_to_managers,
                "teams": [ar.team_id for ar in pres.access_rules if ar.team_id],
                "users": [ar.user_id for ar in pres.access_rules if ar.user_id],
                "type": "html"
            }
            with open(os.path.join(pres_dir, "details.json"), "w") as f:
                json.dump(details, f)

            manifest.append({
                "id": pres.id,
                "title": pres.title,
                "created_at": pres.created_at.strftime("%Y-%m-%d %H:%M"),
                "restricted_to_managers": pres.restricted_to_managers,
                "teams": details["teams"],
                "users": details["users"],
                "slide_count": len(pres.slides)
            })

        with open(os.path.join(base_tmp, "manifest.json"), "w") as f:
            json.dump({"presentations": manifest}, f, indent=2)

        zpath = "/tmp/export_selected_backup.zip"
        if os.path.exists(zpath): os.remove(zpath)

        with zipfile.ZipFile(zpath, "w") as z:
            for root, dirs, files in os.walk(base_tmp):
                for file in files:
                    fpath = os.path.join(root, file)
                    z.write(fpath, os.path.relpath(fpath, base_tmp))

        return send_file(zpath, as_attachment=True, download_name="selected_presentations_backup.zip")

    # GET: show export selection page
    all_pres = Presentation.query.all()
    return render_template("presentations/export_select.html", presentations=all_pres)

@presentations_bp.route("/import_pdf", methods=["GET","POST"])
@login_required
def import_pdf():
    if current_user.role not in ("MANAGER","ADMIN","SUPER_ADMIN"):
        flash("Unauthorized","danger")
        return redirect(url_for("presentations.index"))

    if request.method == "POST":
        title = request.form["title"]
        pdf_file = request.files.get("pdf")
        restricted = bool(request.form.get("restricted_to_managers"))
        reverse = bool(request.form.get("reverse"))
        skip_notes = bool(request.form.get("skip_notes"))

        if not pdf_file or not pdf_file.filename.endswith(".pdf"):
            flash("Please upload a valid PDF file","danger")
            return redirect(url_for("presentations.import_pdf"))

        fname = secure_filename(pdf_file.filename)
        fpath = os.path.join(UPLOAD_MEDIA_DIR, fname)
        pdf_file.save(fpath)

        # âœ… Create new presentation
        pres = Presentation(
            title=title,
            creator_id=current_user.id,
            restricted_to_managers=restricted
        )
        db.session.add(pres)
        db.session.commit()

        # âœ… Convert pages into images
        poppler_path = ensure_poppler()  # Windows returns path, Linux/mac just None
        pages = convert_from_path(fpath, dpi=150, poppler_path=poppler_path)
        if reverse:
            pages = list(reversed(pages))

        for idx, page in enumerate(pages):
            slide_img = f"slide_{pres.id}_{idx}.jpg"
            slide_path = os.path.join(UPLOAD_MEDIA_DIR, "slides", slide_img)
            page.save(slide_path, "JPEG")


            media = MediaFile(filename="slides/"+slide_img, user_id=current_user.id)
            db.session.add(media)


            client_content = f'<img src="/presentations/media/slides/{slide_img}">'
            agent_notes = "" if skip_notes else "(add notes here)"

            db.session.add(Slide(
                presentation_id=pres.id,
                position=idx,
                client_content=client_content,
                agent_notes=agent_notes
            ))

        db.session.commit()
        flash("PDF imported successfully","success")
        return redirect(url_for("presentations.index"))

    return render_template("presentations/import_pdf.html")

